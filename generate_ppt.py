import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # Define some brand colors
    DEEP_FOREST = RGBColor(28, 58, 44)
    CORAL = RGBColor(232, 87, 58)
    INK = RGBColor(74, 74, 63)
    
    # -------------------------------------------------------------
    # Slide 1: Title
    # -------------------------------------------------------------
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Irerero"
    title.text_frame.paragraphs[0].font.color.rgb = DEEP_FOREST
    title.text_frame.paragraphs[0].font.size = Pt(54)
    title.text_frame.paragraphs[0].font.bold = True

    subtitle.text = "Early Childhood Development Platform\nOffline-First Nutritional Monitoring & Management\n\nTeam: [Insert Names Here]"
    for p in subtitle.text_frame.paragraphs:
        p.font.color.rgb = INK

    # -------------------------------------------------------------
    # Slide 2: Agenda
    # -------------------------------------------------------------
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = "Agenda"
    title_shape.text_frame.paragraphs[0].font.color.rgb = DEEP_FOREST

    tf = body_shape.text_frame
    tf.text = "1. Introduction: The Malnutrition Problem in Rwanda"
    p = tf.add_paragraph()
    p.text = "2. System Architecture: How Irerero is Built"
    p = tf.add_paragraph()
    p.text = "3. Project Status: What is Done vs. What is Not Done"
    p = tf.add_paragraph()
    p.text = "4. Live Demo: Seeing Irerero in Action"

    for p in tf.paragraphs:
        p.font.size = Pt(28)
        p.font.color.rgb = INK

    # -------------------------------------------------------------
    # Slide 3: Architecture
    # -------------------------------------------------------------
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    title_shape.text = "System Architecture"
    title_shape.text_frame.paragraphs[0].font.color.rgb = DEEP_FOREST

    tf = shapes.placeholders[1].text_frame
    tf.text = "A Full-Stack, Offline-First Ecosystem:"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(28)

    p = tf.add_paragraph()
    p.text = "Mobile App (Caregivers): Flutter, SQLite for offline storage, Background Sync Engine."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Web Dashboard (Managers/Admins): React, Real-time map & charts, PDF Reporting."
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Backend API: Django REST Framework (Python), PostgreSQL, WHO Z-Score calculation."
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Background Processing: Celery & Redis for asynchronous tasks (Alerts, Reports)."
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Deployment: Hosted on Render cloud infrastructure."
    p.level = 1

    for p in tf.paragraphs:
        p.font.color.rgb = INK

    # -------------------------------------------------------------
    # Slide 4: What is Done vs. What is Not Done
    # -------------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[5]) # Title only layout
    title_shape = slide.shapes.title
    title_shape.text = "Project Status"
    title_shape.text_frame.paragraphs[0].font.color.rgb = DEEP_FOREST

    # Add left box (Done)
    left = Inches(0.5)
    top = Inches(1.8)
    width = Inches(4.2)
    height = Inches(4.5)
    txBox1 = slide.shapes.add_textbox(left, top, width, height)
    tf1 = txBox1.text_frame
    tf1.word_wrap = True
    p = tf1.add_paragraph()
    p.text = "✅ What is Done (Completed):"
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = DEEP_FOREST

    items_done = [
        "Offline-first mobile data entry (attendance & metrics)",
        "Automated nutritional status (SAM/MAM) via WHO standards",
        "Automated Alert & Referral system",
        "Web dashboard with charts & PDF reports",
        "Secure Role-Based Access Control"
    ]
    for item in items_done:
        p = tf1.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(18)
        p.font.color.rgb = INK

    # Add right box (Not Done)
    left = Inches(5.0)
    txBox2 = slide.shapes.add_textbox(left, top, width, height)
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p = tf2.add_paragraph()
    p.text = "🚧 Future Enhancements:"
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = CORAL

    items_not_done = [
        "Hardware integrations (Bluetooth scales/height boards)",
        "Advanced Predictive AI for malnutrition forecasting",
        "Full-scale automated SMS broadcasting gateway for parents"
    ]
    for item in items_not_done:
        p = tf2.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(18)
        p.font.color.rgb = INK

    # -------------------------------------------------------------
    # Slide 5: Live Demo
    # -------------------------------------------------------------
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    title_shape.text = "Live Demo"
    title_shape.text_frame.paragraphs[0].font.color.rgb = DEEP_FOREST

    tf = shapes.placeholders[1].text_frame
    tf.text = "Demonstration Scenarios:"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(28)

    p = tf.add_paragraph()
    p.text = "1. Mobile Offline Data Entry: Caregiver logs attendance and weight/height (triggers SAM/MAM alert)."
    p.level = 1

    p = tf.add_paragraph()
    p.text = "2. Synchronization: Device reconnects to internet and securely pushes data to the server."
    p.level = 1

    p = tf.add_paragraph()
    p.text = "3. Web Dashboard Monitoring: Centre Manager views newly synced data, charts, and reports."
    p.level = 1

    for p in tf.paragraphs:
        p.font.color.rgb = INK

    output_path = r'C:\Users\USER\Desktop\Irerero_Presentation.pptx'
    prs.save(output_path)
    print(f"Presentation successfully saved to: {output_path}")

if __name__ == '__main__':
    create_presentation()
